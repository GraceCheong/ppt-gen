import copy
import datetime
import logging
import os
import shutil
import subprocess
import sys
import tempfile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

from ppt_builder import set_editable_text
from powerpoint_com import create_powerpoint_application, open_presentation_hidden, quit_powerpoint

logger = logging.getLogger(__name__)

_TITLE_FONT = "Diphylleia"
_TITLE_FONT_SIZE = Pt(32)

# 52 natural text colors, one per ISO week of the year.
# Ordered to follow a seasonal progression (winter → spring → summer → autumn → winter).
_WEEK_HEX_COLORS = [
    # Winter (weeks 1–8)
    "#5C6E82", "#6B7280", "#5E6E78", "#667289", "#6E6A82", "#587076", "#6A6E7A", "#5E6A72",
    # Early Spring (weeks 9–17)
    "#6A8260", "#7A8A62", "#6E9070", "#7A9470", "#8A946A", "#9A9460", "#A08A58", "#A8825A",
    "#A87868",
    # Spring Bloom (weeks 18–26)
    "#B07278", "#A87080", "#A06E8A", "#8E6E94", "#7A70A0", "#6A78A8", "#6082A8", "#5A8AA4",
    "#5A9096",
    # Summer (weeks 27–35)
    "#5A946A", "#6A9460", "#7A9458", "#909A58", "#A09A52", "#A88C50", "#B07E50", "#B07260",
    "#B06860",
    # Early Autumn (weeks 36–44)
    "#A86060", "#A05868", "#98587A", "#8A6090", "#7A6898", "#6C6E98", "#687088", "#706878",
    "#887068",
    # Late Autumn / Winter (weeks 45–52)
    "#906858", "#987060", "#A07860", "#A08068", "#989070", "#888A70", "#78886A", "#6A7E78",
]


def _hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def get_week_color(date=None):
    """Return (week_number, RGBColor) for the ISO week of the given date (today if None)."""
    if date is None:
        date = datetime.date.today()
    week_num = date.isocalendar()[1]
    r, g, b = _hex_to_rgb(_WEEK_HEX_COLORS[(week_num - 1) % len(_WEEK_HEX_COLORS)])
    return week_num, RGBColor(r, g, b)


def _get_alt_text(shape):
    try:
        return shape._element.nvSpPr.cNvPr.get("descr", "") or ""
    except AttributeError:
        return ""


def _apply_color(shape, rgb_color):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb_color


_THEME_REL = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
)


def _patch_theme_korean_font(prs, font_name):
    """테마의 한글 스크립트 폰트를 font_name으로 교체한다.

    LibreOffice headless는 run-level <a:ea typeface>를 무시하고
    테마의 <a:font script="Hang"> 값을 한글 글자에 사용한다.
    presentation theme의 majorFont/minorFont 양쪽을 모두 패치해
    LibreOffice가 올바른 폰트로 렌더링하도록 강제한다.
    """
    for master in prs.slide_masters:
        try:
            theme_part = master.part.part_related_by(_THEME_REL)
        except Exception:
            continue
        # theme_part is a generic Part (blob only, no _element)
        root = etree.fromstring(theme_part.blob)
        font_scheme = root.find(".//" + qn("a:fontScheme"))
        if font_scheme is None:
            continue
        changed = False
        for section_tag in ("a:majorFont", "a:minorFont"):
            section = font_scheme.find(qn(section_tag))
            if section is None:
                continue
            # <a:ea typeface>
            ea = section.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(section, qn("a:ea"))
            ea.set("typeface", font_name)
            # <a:font script="Hang"> — LibreOffice가 이 값을 한글 폰트로 사용
            for font_elem in section.findall(qn("a:font")):
                if font_elem.get("script") == "Hang":
                    font_elem.set("typeface", font_name)
                    break
            changed = True
        if changed:
            theme_part._blob = etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            )
            logger.debug("[테마 패치] Hang + ea → %r", font_name)


def build_songlist_pptx(template_path, song_titles, output_pptx_path):
    """Fill the songlist card template with song titles and apply the week color."""
    week_num, rgb = get_week_color()

    prs = Presentation(template_path)
    _patch_theme_korean_font(prs, _TITLE_FONT)
    slide = prs.slides[0]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        alt = _get_alt_text(shape)

        # Fill song title slots: alt text "곡 1", "곡 2", ...
        if alt.startswith("곡"):
            rest = alt[1:].strip()
            if rest.isdigit():
                idx = int(rest) - 1
                title = song_titles[idx] if idx < len(song_titles) else ""
                set_editable_text(shape, title)
                total_runs = sum(len(para.runs) for para in shape.text_frame.paragraphs)
                logger.debug(
                    "[폰트설정] shape=%r alt=%r title=%r paragraphs=%d total_runs=%d",
                    shape.name, alt, title,
                    len(shape.text_frame.paragraphs), total_runs,
                )
                if total_runs == 0:
                    logger.warning(
                        "[폰트설정 실패] shape=%r alt=%r: runs가 없어 폰트를 적용할 수 없습니다. "
                        "set_editable_text 이후 텍스트 프레임에 run이 생성되지 않았습니다.",
                        shape.name, alt,
                    )
                for para in shape.text_frame.paragraphs:
                    # <a:endParaRPr>를 복사해 run의 <a:rPr>로 사용한다.
                    # 이렇게 하면 템플릿이 가진 lang, spc, 텍스트 색상(<a:solidFill>),
                    # <a:cs>, <a:sym> 등 모든 서식이 새 run에 그대로 유지되며
                    # LibreOffice가 일부 글자에만 폰트를 적용하는 문제를 방지한다.
                    endParaRPr = para._p.find(qn("a:endParaRPr"))
                    for run in para.runs:
                        try:
                            if endParaRPr is not None:
                                rPr = copy.deepcopy(endParaRPr)
                                rPr.tag = qn("a:rPr")
                                # latin/ea/cs/sym 모두 _TITLE_FONT로 확인·보정
                                for font_tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
                                    elem = rPr.find(qn(font_tag))
                                    if elem is None:
                                        elem = etree.SubElement(rPr, qn(font_tag))
                                    elem.set("typeface", _TITLE_FONT)
                                existing_rPr = run._r.find(qn("a:rPr"))
                                if existing_rPr is not None:
                                    run._r.remove(existing_rPr)
                                run._r.insert(0, rPr)
                            else:
                                # fallback: endParaRPr가 없는 경우 수동 설정
                                run.font.name = _TITLE_FONT
                                run.font.size = _TITLE_FONT_SIZE
                                rPr = run._r.get_or_add_rPr()
                                for font_tag in ("a:ea", "a:cs", "a:sym"):
                                    elem = rPr.find(qn(font_tag))
                                    if elem is None:
                                        elem = etree.SubElement(rPr, qn(font_tag))
                                    elem.set("typeface", _TITLE_FONT)
                            logger.debug(
                                "[폰트설정 성공] shape=%r run text=%r",
                                shape.name, run.text,
                            )
                        except Exception as e:
                            logger.error(
                                "[폰트설정 오류] shape=%r run text=%r: %s",
                                shape.name, run.text, e,
                                exc_info=True,
                            )

        # Apply week color to "song" and "list" shapes
        if alt.lower() in ("song", "list"):
            _apply_color(shape, rgb)

    prs.save(output_pptx_path)
    return week_num


def _slide_px(pptx_path, long_edge_px=2000):
    prs = Presentation(pptx_path)
    w_emu, h_emu = prs.slide_width, prs.slide_height
    if w_emu >= h_emu:
        return long_edge_px, int(long_edge_px * h_emu / w_emu)
    return int(long_edge_px * w_emu / h_emu), long_edge_px


def _lo_user_install_arg():
    """Return --env:UserInstallation arg pointing to a service-accessible profile dir.

    When running as a Windows Service (SYSTEM account), the default LibreOffice
    user profile path is not writable, so LibreOffice cannot build a font cache
    and silently falls back to built-in fonts even when the font is installed
    system-wide.  Pointing to a ProgramData path fixes this.
    """
    if sys.platform == "win32":
        base = os.environ.get("PROGRAMDATA", r"C:\ProgramData")
        profile_dir = os.path.join(base, "ppt-gen", "lo-profile")
    else:
        profile_dir = os.path.join(os.path.expanduser("~"), ".config", "ppt-gen", "lo-profile")
    os.makedirs(profile_dir, exist_ok=True)
    uri = "file:///" + profile_dir.replace("\\", "/")
    return f"-env:UserInstallation={uri}"


def _lo_subprocess_kwargs():
    """subprocess.run에 전달할 공통 kwargs.

    Windows에서 CREATE_NO_WINDOW 플래그를 설정해 LibreOffice가 콘솔 창을
    상속받지 않도록 한다.  콘솔 창이 없으면 '--version' 등의 출력 후
    나타나는 'Press Enter to continue...' 프롬프트가 차단되지 않는다.
    stdin=DEVNULL 로 LibreOffice가 stdin을 읽으려 해도 즉시 EOF를 받게 한다.
    """
    kwargs = {"stdin": subprocess.DEVNULL, "capture_output": True}
    if sys.platform == "win32":
        kwargs["creationflags"] = subprocess.CREATE_NO_WINDOW
    return kwargs


def find_libreoffice():
    """Return path to LibreOffice executable, or None if not found."""
    if sys.platform == "win32":
        candidates = [
            shutil.which("soffice"),
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    elif sys.platform == "darwin":
        candidates = [
            shutil.which("soffice"),
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
    else:
        candidates = [shutil.which("soffice"), shutil.which("libreoffice")]

    return next((c for c in candidates if c and os.path.isfile(c)), None)


def _export_via_libreoffice(pptx_path, png_path):
    lo = find_libreoffice()
    if not lo:
        return False

    pptx_abs = os.path.abspath(pptx_path)
    png_abs = os.path.abspath(png_path)

    with tempfile.TemporaryDirectory() as tmp:
        # Step 1: PPTX → PDF (LibreOffice PDF 변환이 PNG 직접 변환보다 안정적)
        result = subprocess.run(
            [lo, "--headless", "--norestore", _lo_user_install_arg(),
             "--convert-to", "pdf", "--outdir", tmp, pptx_abs],
            timeout=60,
            **_lo_subprocess_kwargs(),
        )
        if result.returncode != 0:
            logger.warning(
                "[LibreOffice PDF 변환 실패] returncode=%d stderr=%s",
                result.returncode,
                result.stderr.decode("utf-8", errors="replace").strip(),
            )
            return False

        pdf_path = os.path.join(tmp, os.path.splitext(os.path.basename(pptx_abs))[0] + ".pdf")
        if not os.path.exists(pdf_path):
            logger.warning("[LibreOffice PDF 변환 실패] 출력 파일을 찾을 수 없습니다: %s", pdf_path)
            return False

        # Step 2: PDF -> PNG (PyMuPDF)
        try:
            import fitz  # pymupdf

            doc = fitz.open(pdf_path)
            try:
                page = doc[0]
                rect = page.rect
                long_edge = max(rect.width, rect.height)
                scale = 2000 / long_edge if long_edge > 0 else 1.0
                mat = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                pix.save(png_abs)
            finally:
                doc.close()
        except Exception as e:
            logger.warning("[MuPDF PNG 변환 실패] %s", e)
            return False

    return os.path.exists(png_abs)


def _export_via_com(pptx_path, png_path, long_edge_px=None):
    import comtypes
    import comtypes.client

    pptx_abs = os.path.abspath(pptx_path)
    png_abs = os.path.abspath(png_path)
    last_error = None

    long_edges = (long_edge_px,) if long_edge_px else (2000, 1600, 1280)
    for export_long_edge_px in long_edges:
        width_px, height_px = _slide_px(pptx_path, long_edge_px=export_long_edge_px)

        comtypes.CoInitialize()
        try:
            powerpoint = create_powerpoint_application(comtypes.client)
            try:
                prs_com = open_presentation_hidden(powerpoint, pptx_abs)
                try:
                    prs_com.Slides(1).Export(png_abs, "PNG", width_px, height_px)
                    if os.path.exists(png_abs):
                        return
                finally:
                    prs_com.Close()
            finally:
                quit_powerpoint(powerpoint)
        except Exception as e:
            last_error = e
        finally:
            comtypes.CoUninitialize()

    if last_error is not None:
        raise last_error
    raise RuntimeError("PowerPoint COM PNG 변환에 실패했습니다.")


def _export_via_server(pptx_path, png_path, server_url):
    import requests

    url = server_url.rstrip("/") + "/convert"
    with open(pptx_path, "rb") as f:
        response = requests.post(
            url,
            files={"file": ("input.pptx", f, "application/octet-stream")},
            timeout=30,
        )
    response.raise_for_status()

    with open(png_path, "wb") as f:
        f.write(response.content)


def export_pptx_to_png(pptx_path, png_path, server_url=None, long_edge_px=None, skip_com=False):
    """Export the first slide to PNG.

    Priority:
      1. Conversion server (if server_url provided)
      2. Local PowerPoint COM (Windows only, unless skip_com=True)
      3. LibreOffice (if installed locally)
    """
    errors = []

    if server_url:
        try:
            _export_via_server(pptx_path, png_path, server_url)
            logger.info("PNG 변환 완료 [method=server url=%s]", server_url)
            return
        except Exception as e:
            errors.append(f"서버 변환 실패: {e}")

    if sys.platform == "win32" and not skip_com:
        try:
            _export_via_com(pptx_path, png_path, long_edge_px=long_edge_px)
            logger.info("PNG 변환 완료 [method=PowerPoint COM]")
            return
        except Exception as e:
            errors.append(f"로컬 PowerPoint 변환 실패: {e}")

    if _export_via_libreoffice(pptx_path, png_path):
        logger.info("PNG 변환 완료 [method=LibreOffice]")
        return

    errors.append("LibreOffice를 찾을 수 없거나 PNG 변환에 실패했습니다.")

    detail = "\n".join(f"- {message}" for message in errors)
    raise RuntimeError(
        "PPT를 PNG로 변환할 수 없습니다.\n"
        f"{detail}\n\n"
        "서버가 응답하지 않으면 로컬 PowerPoint로 변환을 시도합니다. "
        "로컬 PowerPoint도 사용할 수 없다면 LibreOffice를 설치하세요:\n"
        "https://www.libreoffice.org/download/download-libreoffice/"
    )


def _build_songlist_png_via_com(template_path, song_titles, output_png_path, long_edge_px=None):
    """python-pptx로 PPTX를 수정하고, COM(읽기 전용)으로 PNG를 렌더링한다.

    수정은 python-pptx가 담당하고 COM은 PNG 렌더링만 수행한다.
    read_only=True로 열기 때문에 자동저장/복구 다이얼로그가 발생하지 않는다.
    """
    import comtypes
    import comtypes.client

    # Step 1: python-pptx로 템플릿 수정 → work_pptx에 저장
    _work_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "work")
    os.makedirs(_work_dir, exist_ok=True)
    work_pptx = os.path.join(_work_dir, f"_songlist_{os.getpid()}.pptx")
    week_num = build_songlist_pptx(template_path, song_titles, work_pptx)

    # Step 2: COM으로 PDF로 내보낸 뒤 PDF 1페이지를 PNG로 변환
    png_abs = os.path.abspath(output_png_path)
    long_edges = (long_edge_px,) if long_edge_px else (2000, 1600, 1280)
    last_error = None
    pdf_path = os.path.join(_work_dir, f"_songlist_{os.getpid()}.pdf")

    def _pdf_first_page_to_png(pdf_file, png_file, export_long_edge_px):
        import fitz  # pymupdf

        doc = fitz.open(pdf_file)
        try:
            page = doc[0]
            rect = page.rect
            long_edge = max(rect.width, rect.height)
            scale = export_long_edge_px / long_edge if long_edge > 0 else 1.0
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            pix.save(png_file)
        finally:
            doc.close()

    try:
        for export_long_edge_px in long_edges:
            comtypes.CoInitialize()
            try:
                powerpoint = create_powerpoint_application(comtypes.client)
                try:
                    prs_com = open_presentation_hidden(powerpoint, work_pptx, read_only=True)
                    try:
                        # 32 = ppSaveAsPDF
                        prs_com.SaveAs(pdf_path, 32)
                        _pdf_first_page_to_png(pdf_path, png_abs, export_long_edge_px)
                        if os.path.exists(png_abs):
                            logger.debug("[COM PDF->PNG 내보내기] %r (long_edge=%d)", png_abs, export_long_edge_px)
                            return week_num
                    finally:
                        prs_com.Close()
                finally:
                    quit_powerpoint(powerpoint)
            except Exception as e:
                last_error = e
                logger.warning(
                    "[COM PNG 내보내기 재시도] long_edge=%d error=%s",
                    export_long_edge_px,
                    e,
                )
            finally:
                comtypes.CoUninitialize()
    finally:
        try:
            os.remove(work_pptx)
        except OSError:
            pass
        try:
            os.remove(pdf_path)
        except OSError:
            pass

    if last_error is not None:
        raise last_error
    raise RuntimeError("PowerPoint COM PNG 변환에 실패했습니다.")


def build_songlist_card(template_path, song_titles, output_png_path, server_url=None, skip_com=False):
    """Build the songlist card and export it as PNG. Returns the ISO week number used."""
    # 1순위: COM — 템플릿 수정 + PNG 내보내기를 한 세션에서 처리
    if sys.platform == "win32" and not skip_com:
        try:
            week_num = _build_songlist_png_via_com(template_path, song_titles, output_png_path)
            logger.info("송리스트 카드 생성 완료 [method=COM]")
            return week_num
        except Exception as e:
            logger.warning("COM 변환 실패, python-pptx + LibreOffice로 폴백: %s", e)

    # 2순위 폴백: python-pptx로 PPTX 빌드 → LibreOffice PPTX→PDF→PNG
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        week_num = build_songlist_pptx(template_path, song_titles, tmp_path)
        logger.info("PPTX 생성 완료 [method=python-pptx]")
        export_pptx_to_png(tmp_path, output_png_path, server_url=server_url, skip_com=True)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    return week_num
