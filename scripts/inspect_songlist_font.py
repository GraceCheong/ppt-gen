"""
송리스트 템플릿 및 생성된 PPTX의 폰트 XML을 덤프합니다.

사용법:
    python scripts/inspect_songlist_font.py

출력:
  - 템플릿의 곡 슬롯 shape XML (defRPr, endParaRPr, run rPr)
  - build_songlist_pptx 실행 결과 PPTX의 동일 정보
"""
import os
import sys
import tempfile

ROOT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT_DIR, "src"))

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from songlist_builder import build_songlist_pptx

TEMPLATE_PATH = os.path.join(ROOT_DIR, "assets", "templates", "songlist_template.pptx")
TEST_TITLES = ["믿음이 없이는", "주님 마음 내게 주소서", "주의 자녀로 산다는 것은"]


def _get_alt(shape):
    try:
        return shape._element.nvSpPr.cNvPr.get("descr", "") or ""
    except AttributeError:
        return ""


def dump_songlist_shapes(pptx_path, label):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    print(f"\n{'='*60}")
    print(f"  {label}")
    print(f"  파일: {pptx_path}")
    print(f"{'='*60}")

    found = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        alt = _get_alt(shape)
        if not alt.startswith("곡"):
            continue
        rest = alt[1:].strip()
        if not rest.isdigit():
            continue

        found = True
        print(f"\n[shape={shape.name!r}  alt={alt!r}]")
        txBody = shape.text_frame._txBody
        for para in txBody.findall(qn("a:p")):
            pPr = para.find(qn("a:pPr"))
            if pPr is not None:
                defRPr = pPr.find(qn("a:defRPr"))
                if defRPr is not None:
                    print(f"  <a:defRPr> : {etree.tostring(defRPr, encoding='unicode')}")
                else:
                    print("  <a:defRPr> : (없음)")

            endParaRPr = para.find(qn("a:endParaRPr"))
            if endParaRPr is not None:
                print(f"  <a:endParaRPr>: {etree.tostring(endParaRPr, encoding='unicode')}")

            runs = para.findall(qn("a:r"))
            print(f"  runs 수: {len(runs)}")
            for i, r in enumerate(runs):
                rPr = r.find(qn("a:rPr"))
                t = r.find(qn("a:t"))
                text = t.text if t is not None else ""
                if rPr is not None:
                    print(f"  run[{i}] text={text!r}")
                    print(f"         rPr: {etree.tostring(rPr, encoding='unicode')}")
                else:
                    print(f"  run[{i}] text={text!r}  rPr: (없음 — 상위 defRPr/테마 상속)")

    if not found:
        print("  ⚠ '곡 N' alt text를 가진 shape를 찾지 못했습니다.")


def main():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"[오류] 템플릿 파일을 찾을 수 없습니다: {TEMPLATE_PATH}")
        sys.exit(1)

    # 1. 템플릿 원본 덤프
    dump_songlist_shapes(TEMPLATE_PATH, "템플릿 원본")

    # 2. build_songlist_pptx 실행 후 덤프
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        out_path = tmp.name
    try:
        build_songlist_pptx(TEMPLATE_PATH, TEST_TITLES, out_path)
        dump_songlist_shapes(out_path, f"build_songlist_pptx 결과  (titles={TEST_TITLES})")
    finally:
        try:
            os.remove(out_path)
        except OSError:
            pass


if __name__ == "__main__":
    main()
