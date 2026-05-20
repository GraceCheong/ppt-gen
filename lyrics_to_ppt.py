import re
import os
import sys
import tkinter as tk
import tkinter.font as tkfont
from tkinter import ttk, scrolledtext, messagebox
from pptx import Presentation
try:
    from PIL import Image, ImageTk
except ImportError:
    Image = None
    ImageTk = None

APP_DISPLAY_NAME = "PO,RR"
APP_WINDOW_TITLE = "PO,RR by a tempo"
ASSETS_DIR_NAME = "assets"
ICON_FILE_NAME = "atempo.png"
BACKGROUND_FILE_NAME = "background.png"
BRAND_FONT_CANDIDATES = (
    "Ok Mallang W",
    "나눔스퀘어라운드 ExtraBold",
    "나눔스퀘어 네오 Heavy",
    "나눔스퀘어 네오 ExtraBold",
    "나눔스퀘어 ExtraBold",
    "나눔바른고딕",
    "맑은 고딕",
)

APP_BG = "#eef2f6"
GRADIENT_TOP = "#eef2f6"
GRADIENT_MID = "#a1acbd"
GRADIENT_BOTTOM = "#f4dbe3"
PANEL_BG = "#fffafd"
PANEL_SOFT_BG = "#f1f4f7"
PANEL_BORDER = "#e8c6d0"
TEXT_BG = "#ffffff"
TEXT_FG = "#3d4756"
MUTED_FG = "#6f7a8b"
TITLE_FG = "#3d4756"
ACCENT = "#e7afbf"
ACCENT_DARK = "#bd778c"
ACCENT_SOFT = "#f8e6ec"
LOG_BG = "#3d4756"

LYRICS_GUIDE_TEXT = """V1
가사 첫 번째 줄
가사 두 번째 줄

C
후렴 첫 번째 줄
후렴 두 번째 줄

Out
마지막 가사"""


def hex_to_rgb(color):
    color = color.lstrip("#")
    return tuple(int(color[index:index + 2], 16) for index in (0, 2, 4))


def rgb_to_hex(rgb):
    return "#%02x%02x%02x" % rgb


def blend_hex(start_color, end_color, ratio):
    start_rgb = hex_to_rgb(start_color)
    end_rgb = hex_to_rgb(end_color)
    return rgb_to_hex(
        tuple(
            int(start_rgb[index] + (end_rgb[index] - start_rgb[index]) * ratio)
            for index in range(3)
        )
    )

# ==========================================
# 1. Parsing & Logic Helper Functions 
# ==========================================

def parse_lyrics_text(raw_text):
    lyrics_dict = {}
    blocks = re.split(r'\n\s*\n', raw_text.strip())
    
    for block in blocks:
        lines = block.strip().split('\n')
        if not lines:
            continue
            
        part_key = lines[0].strip()
        lyrics_content = '\n'.join([line.strip() for line in lines[1:]])
        lyrics_dict[part_key] = lyrics_content
        
    return lyrics_dict

def get_base_key(part_key):
    return re.sub(r"'+$", "", part_key)

def chunk_text(text, max_lines=2):
    if not text:
        return []
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    
    MIN_LINE_THRESHOLD = 6
    
    chunks = []
    i = 0
    while i < len(lines):
        take_lines = max_lines
        if max_lines == 1 and len(lines[i]) <= MIN_LINE_THRESHOLD and i + 1 < len(lines):
            take_lines = 2
            
        chunk_lines = lines[i:i+take_lines]
        chunks.append('\n'.join(chunk_lines))
        i += len(chunk_lines)
                
    return chunks

def set_editable_text(shape, text):
    text_frame = shape.text_frame
    text_frame.clear()

    lines = text.split('\n') if text else [""]
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line

def parse_sequence_text(sequence_text):
    lines = [line.strip() for line in sequence_text.splitlines() if line.strip()]

    if not lines:
        raise ValueError("레파토리 입력창이 비어 있습니다.")

    if len(lines) % 2 != 0:
        raise ValueError(f"'{lines[-1]}' 다음에 진행 순서 줄이 없습니다.")

    sequence_entries = []
    for i in range(0, len(lines), 2):
        song_title = lines[i]
        sequence_str = lines[i + 1]

        if not song_title:
            raise ValueError(f"{i + 1}번째 줄의 곡 제목이 비어 있습니다.")
        if not sequence_str:
            raise ValueError(f"'{song_title}'의 진행 순서가 비어 있습니다.")

        sequence_entries.append((song_title, sequence_str))

    return sequence_entries

# ==========================================
# 2. Core PPT Generator Function
# ==========================================

def append_lyrics_to_ppt(prs, song_title, lyrics_text, sequence_str, max_lines_per_slide=2, long_line_threshold=18):
    lyrics_dict = parse_lyrics_text(lyrics_text)
    sequence_list = [part.strip() for part in sequence_str.split('-') if part.strip()]
    # Keep the user's selected line count as a hard maximum. Long lines are
    # handled by the PowerPoint template's text box, not by reducing chunks.
    _ = long_line_threshold
            
    title_layout = None
    lyrics_layout = None
    
    for layout in prs.slide_layouts:
        if "제목" in layout.name:
            title_layout = layout
        if "가사" in layout.name:
            lyrics_layout = layout
            
    if title_layout is None:
        title_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    if lyrics_layout is None:
        lyrics_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[0]

    title_slide = prs.slides.add_slide(title_layout)
    for shape in title_slide.placeholders:
        if shape.has_text_frame:
            set_editable_text(shape, song_title)
            break

    for idx, part in enumerate(sequence_list):
        base_part = get_base_key(part)
        
        if part in lyrics_dict:
            display_text = lyrics_dict[part]
        elif base_part in lyrics_dict:
            display_text = lyrics_dict[base_part]
        else:
            if part.startswith('(') and part.endswith(')'):
                display_text = part[1:-1].strip()
            else:
                if idx == 0 and base_part.upper() in ["I", "INTRO"]:
                    continue
                display_text = "-"
            
        chunks = chunk_text(display_text, max_lines_per_slide)
        if not chunks:
            chunks = [""]
            
        for chunk in chunks:
            slide = prs.slides.add_slide(lyrics_layout)
            
            placeholders = [shape for shape in slide.placeholders if shape.has_text_frame]
            placeholders.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            
            lyrics_placeholder = placeholders[0] if len(placeholders) > 0 else None
            song_title_placeholder = placeholders[1] if len(placeholders) > 1 else None
                            
            if lyrics_placeholder is not None:
                set_editable_text(lyrics_placeholder, chunk)
                
            if song_title_placeholder is not None:
                set_editable_text(song_title_placeholder, song_title)

# ==========================================
# 3. GUI Implementation
# ==========================================

class PillButton(tk.Canvas):
    def __init__(
        self,
        parent,
        text,
        command,
        width=168,
        height=50,
        bg_color="#ffffff",
        hover_bg="#f4f9ff",
        active_bg="#eaf4ff",
        disabled_bg="#eadfea",
        text_color=ACCENT,
        disabled_text="#ad9eb1",
        border_color=PANEL_BORDER,
        font=("Segoe UI", 12, "bold"),
    ):
        super().__init__(
            parent,
            width=width,
            height=height,
            bg=parent.cget("bg"),
            bd=0,
            highlightthickness=0,
            cursor="hand2",
        )
        self.command = command
        self.button_text = text
        self.button_width = width
        self.button_height = height
        self.normal_bg = bg_color
        self.hover_bg = hover_bg
        self.active_bg = active_bg
        self.disabled_bg = disabled_bg
        self.text_color = text_color
        self.disabled_text = disabled_text
        self.border_color = border_color
        self.button_font = font
        self.button_state = tk.NORMAL
        self.current_bg = bg_color

        self.bind("<Enter>", self.on_enter)
        self.bind("<Leave>", self.on_leave)
        self.bind("<ButtonPress-1>", self.on_press)
        self.bind("<ButtonRelease-1>", self.on_release)
        self.draw()

    def rounded_rect(self, x1, y1, x2, y2, radius, **kwargs):
        points = [
            x1 + radius, y1,
            x2 - radius, y1,
            x2, y1,
            x2, y1 + radius,
            x2, y2 - radius,
            x2, y2,
            x2 - radius, y2,
            x1 + radius, y2,
            x1, y2,
            x1, y2 - radius,
            x1, y1 + radius,
            x1, y1,
        ]
        return self.create_polygon(points, smooth=True, **kwargs)

    def draw(self):
        self.delete("all")
        fill = self.disabled_bg if self.button_state == tk.DISABLED else self.current_bg
        foreground = self.disabled_text if self.button_state == tk.DISABLED else self.text_color

        self.rounded_rect(
            2,
            2,
            self.button_width - 2,
            self.button_height - 2,
            self.button_height // 2,
            fill=fill,
            outline=self.border_color,
        )
        self.create_text(
            self.button_width // 2,
            self.button_height // 2,
            text=self.button_text,
            fill=foreground,
            font=self.button_font,
        )

    def on_enter(self, event=None):
        if self.button_state == tk.DISABLED:
            return
        self.current_bg = self.hover_bg
        self.draw()

    def on_leave(self, event=None):
        if self.button_state == tk.DISABLED:
            return
        self.current_bg = self.normal_bg
        self.draw()

    def on_press(self, event=None):
        if self.button_state == tk.DISABLED:
            return
        self.current_bg = self.active_bg
        self.draw()

    def on_release(self, event=None):
        if self.button_state == tk.DISABLED:
            return
        self.current_bg = self.hover_bg
        self.draw()
        if self.command:
            self.command()

    def configure(self, cnf=None, **kwargs):
        if cnf:
            kwargs.update(cnf)

        state = kwargs.pop("state", None)
        text = kwargs.pop("text", None)
        if kwargs:
            super().configure(**kwargs)

        if text is not None:
            self.button_text = text
        if state is not None:
            self.button_state = state
            self.configure(cursor="" if state == tk.DISABLED else "hand2")
        self.draw()

    config = configure


class MultilineDialog(tk.Toplevel):
    def __init__(self, parent, title, prompt):
        super().__init__(parent)
        self.title(title)
        self.geometry("460x380")
        self.transient(parent)
        self.configure(bg=APP_BG)
        self.result = None
        
        content_frame = ttk.Frame(self, padding=(14, 14), style="App.TFrame")
        content_frame.pack(fill=tk.BOTH, expand=True)
        content_frame.rowconfigure(1, weight=1)
        content_frame.columnconfigure(0, weight=1)

        ttk.Label(content_frame, text=prompt, style="Subtle.TLabel").grid(row=0, column=0, sticky=tk.W, pady=(0, 8))
        self.text_area = scrolledtext.ScrolledText(
            content_frame,
            width=48,
            height=12,
            bg=TEXT_BG,
            fg=TEXT_FG,
            insertbackground=TEXT_FG,
            relief=tk.FLAT,
            wrap=tk.WORD,
            font=("맑은 고딕", 10),
            padx=8,
            pady=8,
            highlightthickness=1,
            highlightbackground=PANEL_BORDER,
            highlightcolor=ACCENT,
        )
        self.text_area.grid(row=1, column=0, sticky=tk.NSEW)
        
        btn_frame = ttk.Frame(content_frame, style="App.TFrame")
        btn_frame.grid(row=2, column=0, sticky=tk.E, pady=(12, 0))
        ttk.Button(btn_frame, text="확인", command=self.on_ok, style="Accent.TButton").pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(btn_frame, text="취소", command=self.on_cancel).pack(side=tk.LEFT)
        
        # Enable closing via X button
        self.protocol("WM_DELETE_WINDOW", self.on_cancel)
        
        # Wait for the dialog to be closed
        self.grab_set()
        self.wait_window(self)

    def on_ok(self):
        self.result = self.text_area.get("1.0", tk.END).strip()
        self.destroy()

    def on_cancel(self):
        self.destroy()

class LyricsApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title(APP_WINDOW_TITLE)
        self.geometry("1040x760")
        self.minsize(900, 640)
        
        if getattr(sys, 'frozen', False):
            self.base_dir = os.path.dirname(sys.executable)
        else:
            self.base_dir = os.path.dirname(os.path.abspath(__file__))

        self.configure_window_icon()

        self.sequence_entries = []
        self.current_song_title = None
        self.lyrics_placeholder_visible = False
        self.is_lyrics_dirty = False
        self.loading_lyrics = False
        self.suppress_song_select = False
        self.auto_save_after_id = None
        self.current_song_var = tk.StringVar(value="곡을 선택하세요")

        self.configure(bg=APP_BG)
        self._background_image = None
        self._background_cache_key = None
        self.brand_font_family = self.resolve_font_family(BRAND_FONT_CANDIDATES)
        self.setup_style()
        self.create_widgets()
        self.load_initial_repertoire_text()
        self.protocol("WM_DELETE_WINDOW", self.on_close)

    def configure_window_icon(self):
        icon_file = self.find_asset_file(ICON_FILE_NAME)
        if icon_file:
            try:
                self._window_icon = tk.PhotoImage(file=icon_file)
                self.iconphoto(True, self._window_icon)
                return
            except tk.TclError:
                pass

        try:
            self._empty_icon = tk.PhotoImage(width=1, height=1)
            self.iconphoto(True, self._empty_icon)
        except tk.TclError:
            try:
                self.iconbitmap(default="")
            except tk.TclError:
                pass

    def find_icon_file(self):
        return self.find_asset_file(ICON_FILE_NAME)

    def find_background_file(self):
        return self.find_asset_file(BACKGROUND_FILE_NAME)

    def find_asset_file(self, file_name):
        base_dirs = [self.base_dir]
        bundle_dir = getattr(sys, "_MEIPASS", None)
        if bundle_dir:
            base_dirs.append(bundle_dir)

        relative_paths = [
            os.path.join(ASSETS_DIR_NAME, file_name),
            file_name,
        ]

        for base_dir in base_dirs:
            for relative_path in relative_paths:
                asset_file = os.path.join(base_dir, relative_path)
                if os.path.exists(asset_file):
                    return asset_file

        return None

    def resolve_font_family(self, candidates):
        available_fonts = set(tkfont.families(self))
        for candidate in candidates:
            if candidate in available_fonts:
                return candidate
        return "맑은 고딕"

    def setup_style(self):
        style = ttk.Style(self)
        try:
            style.theme_use("clam")
        except tk.TclError:
            pass

        style.configure(".", font=("Segoe UI", 10))
        style.configure("App.TFrame", background=APP_BG)
        style.configure("Panel.TFrame", background=PANEL_BG)
        style.configure("Title.TLabel", background=APP_BG, foreground=TITLE_FG, font=("Segoe UI", 30, "bold"))
        style.configure("Subtle.TLabel", background=APP_BG, foreground=MUTED_FG, font=("Segoe UI", 10))
        style.configure("Panel.TLabel", background=PANEL_BG, foreground=TEXT_FG, font=("Segoe UI", 10, "bold"))
        style.configure("Field.TLabel", background=PANEL_BG, foreground=TEXT_FG, font=("Segoe UI", 9, "bold"))
        style.configure("Hint.TLabel", background=PANEL_BG, foreground=MUTED_FG, font=("Segoe UI", 9))
        style.configure("TButton", font=("Segoe UI", 10), padding=(12, 7), background="#ffffff")
        style.map("TButton", background=[("active", ACCENT_SOFT), ("disabled", "#e7dce1")])
        style.configure("Accent.TButton", font=("Segoe UI", 10, "bold"), padding=(14, 8), foreground="#ffffff", background=ACCENT)
        style.map("Accent.TButton", background=[("active", ACCENT_DARK), ("disabled", "#9fb6b8")])
        style.configure("TSpinbox", fieldbackground=TEXT_BG, background=TEXT_BG)

    def create_widgets(self):
        self.background_canvas = tk.Canvas(self, bd=0, highlightthickness=0)
        self.background_canvas.pack(fill=tk.BOTH, expand=True)
        self.background_canvas.bind("<Configure>", self.on_canvas_configure)

        self.title_item = self.background_canvas.create_text(
            0,
            0,
            text=APP_DISPLAY_NAME,
            fill=TITLE_FG,
            font=(self.brand_font_family, 38, "bold"),
            tags=("foreground",),
        )
        self.subtitle_item = self.background_canvas.create_text(
            0,
            0,
            text="레파토리와 가사를 정리해 파워포인트로 만듭니다.",
            fill=MUTED_FG,
            font=("Segoe UI", 11, "bold"),
            tags=("foreground",),
        )

        settings_frame = self.create_panel(self.background_canvas)
        self.settings_window = self.background_canvas.create_window(
            0, 0, anchor=tk.NW, window=settings_frame, tags=("foreground",)
        )
        settings_frame.columnconfigure(3, weight=1)

        tk.Label(
            settings_frame,
            text="생성 설정",
            bg=PANEL_BG,
            fg=TEXT_FG,
            font=("Segoe UI", 12, "bold"),
        ).grid(row=0, column=0, sticky=tk.W, padx=(18, 16), pady=16)
        
        ttk.Label(settings_frame, text="슬라이드당 줄 수", style="Field.TLabel").grid(row=0, column=1, sticky=tk.W)
        self.max_lines_var = tk.IntVar(value=2)
        ttk.Spinbox(settings_frame, from_=1, to=10, textvariable=self.max_lines_var, width=5).grid(row=0, column=2, sticky=tk.W, padx=(8, 12))
        self.threshold_var = tk.IntVar(value=18)
        ttk.Label(settings_frame, text="설정한 줄 수를 그대로 사용합니다.", style="Hint.TLabel").grid(row=0, column=3, sticky=tk.W)

        workspace_frame = self.create_panel(self.background_canvas)
        self.workspace_window = self.background_canvas.create_window(
            0, 0, anchor=tk.NW, window=workspace_frame, tags=("foreground",)
        )
        workspace_frame.columnconfigure(0, weight=4)
        workspace_frame.columnconfigure(1, weight=5)
        workspace_frame.rowconfigure(0, weight=1)

        sequence_frame = tk.Frame(workspace_frame, bg=PANEL_BG, bd=0, highlightthickness=0)
        sequence_frame.grid(row=0, column=0, sticky=tk.NSEW, padx=(16, 8), pady=16)
        sequence_frame.rowconfigure(1, weight=1)
        sequence_frame.columnconfigure(0, weight=1)

        tk.Label(
            sequence_frame,
            text="레파토리 입력",
            bg=PANEL_BG,
            fg=TEXT_FG,
            font=("Segoe UI", 13, "bold"),
        ).grid(row=0, column=0, sticky=tk.W, padx=16, pady=(14, 8))

        self.sequence_text = scrolledtext.ScrolledText(
            sequence_frame,
            width=42,
            height=14,
            bg=TEXT_BG,
            fg=TEXT_FG,
            insertbackground=TEXT_FG,
            relief=tk.FLAT,
            wrap=tk.WORD,
            font=("맑은 고딕", 10),
            padx=8,
            pady=8,
            highlightthickness=1,
            highlightbackground=PANEL_BORDER,
            highlightcolor=ACCENT,
        )
        self.sequence_text.grid(row=1, column=0, sticky=tk.NSEW, padx=16, pady=(0, 16))

        lyrics_frame = tk.Frame(workspace_frame, bg=PANEL_BG, bd=0, highlightthickness=0)
        lyrics_frame.grid(row=0, column=1, sticky=tk.NSEW, padx=(8, 16), pady=16)
        lyrics_frame.columnconfigure(0, weight=0)
        lyrics_frame.columnconfigure(1, weight=1)
        lyrics_frame.rowconfigure(1, weight=1)

        tk.Label(
            lyrics_frame,
            text="가사 편집",
            bg=PANEL_BG,
            fg=TEXT_FG,
            font=("Segoe UI", 13, "bold"),
        ).grid(row=0, column=0, sticky=tk.W, padx=16, pady=(14, 8))
        tk.Label(
            lyrics_frame,
            textvariable=self.current_song_var,
            bg=PANEL_BG,
            fg=MUTED_FG,
            font=("맑은 고딕", 10, "bold"),
        ).grid(row=0, column=1, sticky=tk.W, padx=(0, 16), pady=(14, 8))

        list_frame = tk.Frame(lyrics_frame, bg=PANEL_BG, bd=0, highlightthickness=0)
        list_frame.grid(row=1, column=0, sticky=tk.NS, padx=(16, 10), pady=(0, 16))

        self.song_list = tk.Listbox(
            list_frame,
            width=20,
            height=12,
            bg=PANEL_SOFT_BG,
            fg=TEXT_FG,
            selectbackground=ACCENT,
            selectforeground="#ffffff",
            highlightthickness=1,
            highlightbackground="#d8dee8",
            highlightcolor=ACCENT,
            relief=tk.FLAT,
            activestyle="none",
            font=("맑은 고딕", 10),
        )
        self.song_list.pack(side=tk.LEFT, fill=tk.Y)
        self.song_list.bind("<<ListboxSelect>>", self.on_song_select)

        song_scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.song_list.yview)
        song_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.song_list.config(yscrollcommand=song_scrollbar.set)

        self.lyrics_text = scrolledtext.ScrolledText(
            lyrics_frame,
            width=48,
            height=14,
            bg=TEXT_BG,
            fg=TEXT_FG,
            insertbackground=TEXT_FG,
            relief=tk.FLAT,
            wrap=tk.WORD,
            font=("맑은 고딕", 10),
            padx=8,
            pady=8,
            highlightthickness=1,
            highlightbackground=PANEL_BORDER,
            highlightcolor=ACCENT,
        )
        self.lyrics_text.grid(row=1, column=1, sticky=tk.NSEW, padx=(0, 16), pady=(0, 16))
        self.lyrics_text.tag_configure("placeholder", foreground="#9aa3af")
        self.lyrics_text.bind("<FocusIn>", self.on_lyrics_focus_in)
        self.lyrics_text.bind("<FocusOut>", self.on_lyrics_focus_out)
        self.lyrics_text.bind("<<Modified>>", self.on_lyrics_modified)
        self.show_lyrics_guide()

        action_frame = self.create_panel(self.background_canvas)
        self.action_window = self.background_canvas.create_window(
            0, 0, anchor=tk.NW, window=action_frame, tags=("foreground",)
        )

        self.refresh_btn = PillButton(
            action_frame,
            "레파토리 인식",
            self.refresh_song_list,
            width=170,
            bg_color="#ffffff",
            hover_bg=ACCENT_SOFT,
            active_bg="#f2d4de",
            text_color=TEXT_FG,
            border_color=ACCENT,
        )
        self.refresh_btn.pack(side=tk.LEFT, padx=8, pady=10)

        self.save_btn = PillButton(
            action_frame,
            "가사 저장",
            self.save_current_lyrics,
            width=142,
            bg_color="#ffffff",
            hover_bg=ACCENT_SOFT,
            active_bg="#f2d4de",
            text_color=TEXT_FG,
            border_color=ACCENT,
        )
        self.save_btn.pack(side=tk.LEFT, padx=8, pady=10)

        self.download_btn = PillButton(
            action_frame,
            "가사 다운로드",
            self.download_lyrics,
            width=158,
            bg_color="#ffffff",
            hover_bg=ACCENT_SOFT,
            active_bg="#f2d4de",
            text_color=TEXT_FG,
            border_color=ACCENT,
        )
        self.download_btn.pack(side=tk.LEFT, padx=8, pady=10)

        self.generate_btn = PillButton(
            action_frame,
            "파워포인트 생성",
            self.generate_ppt,
            width=186,
            bg_color=ACCENT,
            hover_bg="#efc4d0",
            active_bg="#d99bab",
            text_color=TEXT_FG,
            border_color=ACCENT_DARK,
            font=("Segoe UI", 13, "bold"),
        )
        self.generate_btn.pack(side=tk.RIGHT, padx=8, pady=10)

        log_frame = self.create_panel(self.background_canvas)
        self.log_window = self.background_canvas.create_window(
            0, 0, anchor=tk.NW, window=log_frame, tags=("foreground",)
        )
        log_frame.rowconfigure(1, weight=1)
        log_frame.columnconfigure(0, weight=1)

        tk.Label(
            log_frame,
            text="작업 로그",
            bg=PANEL_BG,
            fg=TEXT_FG,
            font=("Segoe UI", 12, "bold"),
        ).grid(row=0, column=0, sticky=tk.W, padx=16, pady=(12, 8))

        self.log_area = scrolledtext.ScrolledText(
            log_frame,
            width=70,
            height=9,
            state=tk.DISABLED,
            bg=LOG_BG,
            fg="#f6edf1",
            insertbackground="#f6edf1",
            relief=tk.FLAT,
            font=("Consolas", 9),
        )
        self.log_area.grid(row=1, column=0, sticky=tk.NSEW, padx=16, pady=(0, 16))

    def create_panel(self, parent):
        return tk.Frame(
            parent,
            bg=PANEL_BG,
            bd=0,
            highlightthickness=1,
            highlightbackground=PANEL_BORDER,
            highlightcolor=PANEL_BORDER,
        )

    def draw_background(self, width, height):
        self.background_canvas.delete("background")
        background_file = self.find_background_file()

        if background_file:
            self.draw_background_image(width, height, background_file)
            self.background_canvas.tag_lower("background")
            return

        split = max(1, int(height * 0.55))

        for y in range(height):
            if y < split:
                color = blend_hex(GRADIENT_TOP, GRADIENT_MID, y / split)
            else:
                color = blend_hex(GRADIENT_MID, GRADIENT_BOTTOM, (y - split) / max(1, height - split))
            self.background_canvas.create_line(0, y, width, y, fill=color, tags=("background",))

        self.background_canvas.create_oval(
            int(width * 0.58),
            int(height * 0.52),
            int(width * 1.18),
            int(height * 1.18),
            fill="#eadcf4",
            outline="",
            tags=("background",),
        )
        self.background_canvas.create_oval(
            -int(width * 0.20),
            int(height * 0.36),
            int(width * 0.38),
            int(height * 1.12),
            fill="#dfcaef",
            outline="",
            tags=("background",),
        )
        self.background_canvas.tag_lower("background")

    def draw_background_image(self, width, height, background_file):
        cache_key = (background_file, width, height)

        if Image is not None and ImageTk is not None:
            if self._background_cache_key != cache_key:
                with Image.open(background_file) as source_image:
                    source_image = source_image.convert("RGB")
                    image_ratio = source_image.width / source_image.height
                    target_ratio = width / max(1, height)

                    if image_ratio > target_ratio:
                        new_height = height
                        new_width = max(width, int(height * image_ratio))
                    else:
                        new_width = width
                        new_height = max(height, int(width / image_ratio))

                    resized = source_image.resize((new_width, new_height), Image.LANCZOS)
                    left = max(0, (new_width - width) // 2)
                    top = max(0, (new_height - height) // 2)
                    cropped = resized.crop((left, top, left + width, top + height))

                    palette_overlay = Image.new("RGB", (1, height), PANEL_BG)
                    palette_overlay.putdata(
                        [
                            hex_to_rgb(blend_hex(GRADIENT_TOP, GRADIENT_BOTTOM, y / max(1, height - 1)))
                            for y in range(height)
                        ]
                    )
                    palette_overlay = palette_overlay.resize((width, height))

                    toned = Image.blend(cropped, palette_overlay, 0.38)
                    self._background_image = ImageTk.PhotoImage(toned)
                    self._background_cache_key = cache_key

            self.background_canvas.create_image(
                0,
                0,
                anchor=tk.NW,
                image=self._background_image,
                tags=("background",),
            )
            return

        try:
            self._background_image = tk.PhotoImage(file=background_file)
            self._background_cache_key = cache_key
            self.background_canvas.create_image(
                width // 2,
                height // 2,
                anchor=tk.CENTER,
                image=self._background_image,
                tags=("background",),
            )
        except tk.TclError:
            for y in range(height):
                color = blend_hex(GRADIENT_TOP, GRADIENT_BOTTOM, y / max(1, height))
                self.background_canvas.create_line(0, y, width, y, fill=color, tags=("background",))

    def on_canvas_configure(self, event):
        width = event.width
        height = event.height
        self.draw_background(width, height)

        margin = 28
        gap = 12
        header_h = 104
        settings_h = 74
        action_h = 72
        log_h = max(118, min(150, int(height * 0.18)))
        content_w = max(1, width - margin * 2)
        workspace_h = max(
            260,
            height - margin * 2 - header_h - settings_h - action_h - log_h - gap * 4,
        )

        settings_y = margin + header_h
        workspace_y = settings_y + settings_h + gap
        action_y = workspace_y + workspace_h + gap
        log_y = action_y + action_h + gap

        self.background_canvas.coords(self.title_item, width // 2, margin + 32)
        self.background_canvas.coords(self.subtitle_item, width // 2, margin + 78)

        self.background_canvas.coords(self.settings_window, margin, settings_y)
        self.background_canvas.itemconfigure(self.settings_window, width=content_w, height=settings_h)

        self.background_canvas.coords(self.workspace_window, margin, workspace_y)
        self.background_canvas.itemconfigure(self.workspace_window, width=content_w, height=workspace_h)

        self.background_canvas.coords(self.action_window, margin, action_y)
        self.background_canvas.itemconfigure(self.action_window, width=content_w, height=action_h)

        self.background_canvas.coords(self.log_window, margin, log_y)
        self.background_canvas.itemconfigure(self.log_window, width=content_w, height=log_h)
        
    def log(self, message):
        self.log_area.config(state=tk.NORMAL)
        self.log_area.insert(tk.END, message + "\n")
        self.log_area.see(tk.END)
        self.log_area.config(state=tk.DISABLED)
        self.update_idletasks()

    def get_manual_lyrics(self, song_title):
        dialog = MultilineDialog(self, "가사 직접 입력", f"'{song_title}' 가사를 입력하세요.")
        return dialog.result or ""

    def get_output_dir(self):
        output_dir = os.path.join(self.base_dir, "out")
        os.makedirs(output_dir, exist_ok=True)
        return output_dir

    def get_sequence_entries(self):
        sequence_text = self.sequence_text.get("1.0", tk.END)
        return parse_sequence_text(sequence_text)

    def set_action_buttons_state(self, state):
        self.refresh_btn.config(state=state)
        self.save_btn.config(state=state)
        self.download_btn.config(state=state)
        self.generate_btn.config(state=state)

    def lyrics_file_path(self, song_title):
        return os.path.join(self.base_dir, f"{song_title}.txt")

    def load_initial_repertoire_text(self):
        sequence_file = os.path.join(self.base_dir, "sequences.txt")
        if not os.path.exists(sequence_file):
            return

        try:
            with open(sequence_file, "r", encoding="utf-8-sig") as f:
                repertoire_text = f.read().strip()
        except Exception as e:
            self.log(f"[오류] 기존 레파토리 파일을 읽지 못했습니다: {e}")
            return

        if not repertoire_text:
            return

        self.sequence_text.insert("1.0", repertoire_text)

        try:
            sequence_entries = parse_sequence_text(repertoire_text)
        except ValueError as e:
            self.log(f"[안내] sequences.txt를 불러왔지만 형식 확인이 필요합니다: {e}")
            return

        self.sequence_entries = sequence_entries
        selected_index = self.populate_song_list(sequence_entries, preserve_current=False)
        if selected_index is not None:
            self.load_lyrics_for_song(sequence_entries[selected_index][0])

        self.log(f"[안내] sequences.txt에서 레파토리 {len(sequence_entries)}곡을 불러왔습니다.")

    def show_lyrics_guide(self):
        self.loading_lyrics = True
        self.lyrics_text.config(state=tk.NORMAL)
        self.lyrics_text.delete("1.0", tk.END)
        self.lyrics_text.insert("1.0", LYRICS_GUIDE_TEXT, "placeholder")
        self.lyrics_placeholder_visible = True
        self.is_lyrics_dirty = False
        self.lyrics_text.edit_modified(False)
        self.loading_lyrics = False

    def clear_lyrics_guide(self):
        if not self.lyrics_placeholder_visible:
            return

        self.loading_lyrics = True
        self.lyrics_text.delete("1.0", tk.END)
        self.lyrics_placeholder_visible = False
        self.is_lyrics_dirty = False
        self.lyrics_text.edit_modified(False)
        self.loading_lyrics = False

    def set_lyrics_editor_text(self, text):
        self.cancel_scheduled_auto_save()
        self.loading_lyrics = True
        self.lyrics_text.config(state=tk.NORMAL)
        self.lyrics_text.delete("1.0", tk.END)
        self.lyrics_text.insert("1.0", text)
        self.lyrics_placeholder_visible = False
        self.is_lyrics_dirty = False
        self.lyrics_text.edit_modified(False)
        self.loading_lyrics = False

    def get_lyrics_editor_text(self):
        if self.lyrics_placeholder_visible:
            return ""
        return self.lyrics_text.get("1.0", tk.END).strip()

    def on_lyrics_focus_in(self, event=None):
        self.clear_lyrics_guide()

    def on_lyrics_focus_out(self, event=None):
        if not self.get_lyrics_editor_text():
            self.show_lyrics_guide()

    def on_lyrics_modified(self, event=None):
        if self.loading_lyrics:
            self.lyrics_text.edit_modified(False)
            return

        if self.lyrics_text.edit_modified():
            if self.current_song_title and not self.lyrics_placeholder_visible:
                self.is_lyrics_dirty = True
                self.schedule_lyrics_auto_save()
            self.lyrics_text.edit_modified(False)

    def schedule_lyrics_auto_save(self):
        self.cancel_scheduled_auto_save()
        self.auto_save_after_id = self.after(1200, self.run_scheduled_auto_save)

    def cancel_scheduled_auto_save(self):
        if self.auto_save_after_id is None:
            return

        try:
            self.after_cancel(self.auto_save_after_id)
        except tk.TclError:
            pass
        self.auto_save_after_id = None

    def run_scheduled_auto_save(self):
        self.auto_save_after_id = None
        self.auto_save_current_lyrics()

    def populate_song_list(self, sequence_entries, preserve_current=True):
        previous_song = self.current_song_title if preserve_current else None
        selected_index = None

        self.suppress_song_select = True
        self.song_list.delete(0, tk.END)

        for index, (song_title, _) in enumerate(sequence_entries):
            self.song_list.insert(tk.END, song_title)
            if previous_song == song_title and selected_index is None:
                selected_index = index

        if selected_index is None and sequence_entries:
            selected_index = 0

        if selected_index is not None:
            self.song_list.selection_set(selected_index)
            self.song_list.activate(selected_index)

        self.suppress_song_select = False
        return selected_index

    def restore_song_selection(self):
        self.suppress_song_select = True
        self.song_list.selection_clear(0, tk.END)

        if self.current_song_title:
            for index in range(self.song_list.size()):
                if self.song_list.get(index) == self.current_song_title:
                    self.song_list.selection_set(index)
                    self.song_list.activate(index)
                    break

        self.suppress_song_select = False

    def auto_save_current_lyrics(self):
        if not self.current_song_title or not self.is_lyrics_dirty:
            return True

        return self.save_current_lyrics(
            show_message=False,
            allow_empty=True,
            log_success=False,
        )

    def refresh_song_list(self, show_message=True):
        if not self.auto_save_current_lyrics():
            return False

        try:
            sequence_entries = self.get_sequence_entries()
        except ValueError as e:
            self.log(f"[오류] {e}")
            messagebox.showerror("레파토리 입력 오류", str(e))
            return False

        self.sequence_entries = sequence_entries
        selected_index = self.populate_song_list(sequence_entries)

        if selected_index is not None:
            song_title = sequence_entries[selected_index][0]
            if song_title != self.current_song_title:
                self.load_lyrics_for_song(song_title)
        else:
            self.current_song_title = None
            self.current_song_var.set("곡을 선택하세요")
            self.show_lyrics_guide()

        if show_message:
            self.log(f"[안내] 레파토리 {len(sequence_entries)}곡을 인식했습니다.")

        return True

    def on_song_select(self, event=None):
        if self.suppress_song_select:
            return

        selection = self.song_list.curselection()
        if not selection:
            return

        song_title = self.song_list.get(selection[0])
        if song_title == self.current_song_title:
            return

        if not self.auto_save_current_lyrics():
            self.restore_song_selection()
            return

        self.load_lyrics_for_song(song_title)

    def load_lyrics_for_song(self, song_title):
        self.current_song_title = song_title
        self.current_song_var.set(song_title)

        lyrics_file = self.lyrics_file_path(song_title)
        if os.path.exists(lyrics_file):
            try:
                with open(lyrics_file, "r", encoding="utf-8") as f:
                    raw_lyrics = f.read()
            except Exception as e:
                self.log(f"[오류] '{lyrics_file}' 파일을 읽지 못했습니다: {e}")
                self.show_lyrics_guide()
                return

            if raw_lyrics.strip():
                self.set_lyrics_editor_text(raw_lyrics)
            else:
                self.show_lyrics_guide()
        else:
            self.show_lyrics_guide()

    def save_current_lyrics(self, show_message=True, allow_empty=False, log_success=True):
        if not self.current_song_title:
            messagebox.showwarning("가사 저장", "레파토리 인식 후 곡을 먼저 선택하세요.")
            return False

        lyrics_text = self.get_lyrics_editor_text()
        if not lyrics_text and not allow_empty:
            messagebox.showwarning("가사 저장", "저장할 가사를 입력하세요.")
            return False

        self.cancel_scheduled_auto_save()
        lyrics_file = self.lyrics_file_path(self.current_song_title)
        try:
            with open(lyrics_file, "w", encoding="utf-8") as f:
                f.write(lyrics_text)
        except Exception as e:
            self.log(f"[오류] '{lyrics_file}' 파일을 저장하지 못했습니다: {e}")
            messagebox.showerror("오류", f"가사를 저장하지 못했습니다:\n{e}")
            return False

        self.is_lyrics_dirty = False
        self.lyrics_text.edit_modified(False)
        if log_success:
            self.log(f"[완료] 가사를 저장했습니다: '{lyrics_file}'")

        if show_message:
            messagebox.showinfo("가사 저장", "가사를 저장했습니다.")

        return True

    def on_close(self):
        if self.auto_save_current_lyrics():
            self.destroy()

    def download_lyrics(self):
        if not self.refresh_song_list(show_message=False):
            return

        try:
            from auto_lyrics_downloader import download_missing_lyrics
        except Exception as e:
            self.log(f"[오류] 가사 다운로드 모듈을 불러오지 못했습니다: {e}")
            messagebox.showerror("오류", f"가사 다운로드 모듈을 불러오지 못했습니다:\n{e}")
            return

        sequence_entries = self.sequence_entries
        song_titles = [song_title for song_title, _ in sequence_entries]
        current_song = self.current_song_title

        self.set_action_buttons_state(tk.DISABLED)
        self.log("====================================")
        self.log("가사 다운로드를 시작합니다.")

        try:
            download_missing_lyrics(
                song_titles=song_titles,
                base_dir=self.base_dir,
                log_func=self.log,
            )
            if current_song and not self.is_lyrics_dirty:
                self.load_lyrics_for_song(current_song)
            messagebox.showinfo("완료", "가사 다운로드 작업이 완료되었습니다.")
        except Exception as e:
            self.log(f"[오류] 가사 다운로드에 실패했습니다: {e}")
            messagebox.showerror("오류", f"가사 다운로드에 실패했습니다:\n{e}")
        finally:
            self.set_action_buttons_state(tk.NORMAL)

    def generate_ppt(self):
        self.log("====================================")
        self.log("파워포인트 생성을 시작합니다.")

        if not self.refresh_song_list(show_message=False):
            return

        sequence_entries = self.sequence_entries
        
        max_lines_per_slide = self.max_lines_var.get()
        long_line_threshold = self.threshold_var.get()
        
        template_file = os.path.join(self.base_dir, "template.pptx")
        
        if not os.path.exists(template_file):
            self.log(f"[오류] 템플릿 파일을 찾을 수 없습니다: '{template_file}'")
            messagebox.showerror("오류", f"템플릿 파일을 찾을 수 없습니다:\n{template_file}")
            return
            
        try:
            prs = Presentation(template_file)
        except Exception as e:
            self.log(f"[오류] 템플릿을 불러오지 못했습니다: {e}")
            messagebox.showerror("오류", f"템플릿을 불러오지 못했습니다:\n{e}")
            return

        appended_count = 0
        for song_title, sequence_str in sequence_entries:
            lyrics_file = self.lyrics_file_path(song_title)
            raw_lyrics = ""
            
            if os.path.exists(lyrics_file):
                self.log(f"[진행] '{song_title}' 처리 중")
                try:
                    with open(lyrics_file, 'r', encoding='utf-8') as f:
                        raw_lyrics = f.read()
                except Exception as e:
                    self.log(f"[오류] '{lyrics_file}' 파일을 읽지 못했습니다: {e}")
                    raw_lyrics = self.get_manual_lyrics(song_title)
            else:
                self.log(f"[안내] '{lyrics_file}' 파일이 없어 직접 입력 창을 엽니다.")
                raw_lyrics = self.get_manual_lyrics(song_title)
                
            if raw_lyrics.strip():
                append_lyrics_to_ppt(prs, song_title, raw_lyrics, sequence_str, max_lines_per_slide, long_line_threshold)
                appended_count += 1
            else:
                self.log(f"[안내] '{song_title}' 가사가 없어 건너뜁니다.")

        if appended_count == 0:
            self.log("[오류] 생성할 가사가 없습니다.")
            messagebox.showwarning("파워포인트 생성", "생성할 가사가 없습니다.")
            return

        output_file = os.path.join(self.get_output_dir(), "integrated_lyrics.pptx")
        try:
            prs.save(output_file)
            self.log(f"\n[완료] 파워포인트 파일을 만들었습니다: '{output_file}'\n")
            messagebox.showinfo("완료", "파워포인트 파일을 생성했습니다.\n저장 위치: out/integrated_lyrics.pptx")
        except Exception as e:
            self.log(f"[오류] 파워포인트 파일을 저장하지 못했습니다: {e}")
            messagebox.showerror("오류", f"파워포인트 파일을 저장하지 못했습니다:\n{e}")

if __name__ == "__main__":
    app = LyricsApp()
    app.mainloop()
