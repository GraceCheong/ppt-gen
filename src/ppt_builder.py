import re
from copy import deepcopy

from lxml import etree
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_MIN_SHORT_LINE_LEN = 6


def _a(tag):
    return f'{{{_NS}}}{tag}'


def _as_run_properties(element):
    """Clone DrawingML character properties so they can be used as ``a:rPr``."""
    if element is None:
        return None
    cloned = deepcopy(element)
    cloned.tag = _a('rPr')
    return cloned


def _character_properties_from_text_body(tx_body):
    """Return the most specific character properties stored in a text body.

    PowerPoint commonly stores placeholder fonts in ``endParaRPr`` or
    ``lstStyle/lvlXpPr/defRPr`` rather than in a run. python-pptx does not
    resolve that inheritance when new text is inserted, so inspect those
    locations explicitly.
    """
    candidates = []
    paragraphs = tx_body.findall(_a('p'))
    if paragraphs:
        paragraph = paragraphs[0]
        run = paragraph.find(_a('r'))
        if run is not None:
            run_props = run.find(_a('rPr'))
            if run_props is not None:
                candidates.append(run_props)

        end_props = paragraph.find(_a('endParaRPr'))
        if end_props is not None:
            candidates.append(end_props)

        paragraph_props = paragraph.find(_a('pPr'))
        if paragraph_props is not None:
            default_props = paragraph_props.find(_a('defRPr'))
            if default_props is not None:
                candidates.append(default_props)

    list_style = tx_body.find(_a('lstStyle'))
    if list_style is not None:
        for level in list_style:
            default_props = level.find(_a('defRPr'))
            if default_props is not None:
                candidates.append(default_props)

    if not candidates:
        return None

    # Candidates are ordered most-specific to least-specific. Start with the
    # most-specific element, then fill only missing values from inherited
    # levels. A common PowerPoint pattern is language in endParaRPr and the
    # actual Korean font in lstStyle/defRPr.
    merged = _as_run_properties(candidates[0])
    for fallback in candidates[1:]:
        for name, value in fallback.attrib.items():
            if name not in merged.attrib:
                merged.set(name, value)
        existing_tags = {child.tag for child in merged}
        for child in fallback:
            if child.tag not in existing_tags:
                merged.append(deepcopy(child))
                existing_tags.add(child.tag)
    return merged


def _theme_font_variant(typeface, script):
    if typeface == '+mj-lt':
        return f'+mj-{script}'
    if typeface == '+mn-lt':
        return f'+mn-{script}'
    return typeface


def _ensure_script_fonts(run_props):
    """Give Latin, Korean/East-Asian and complex text explicit font entries."""
    font_elements = {
        tag: run_props.find(_a(tag)) for tag in ('latin', 'ea', 'cs')
    }
    source = next(
        (
            element.get('typeface')
            for element in font_elements.values()
            if element is not None and element.get('typeface')
        ),
        None,
    )
    if not source:
        return run_props

    for tag, script in (('latin', 'lt'), ('ea', 'ea'), ('cs', 'cs')):
        element = font_elements[tag]
        if element is None:
            element = etree.SubElement(run_props, _a(tag))
        if not element.get('typeface'):
            element.set('typeface', _theme_font_variant(source, script))
    return run_props


def _default_theme_run_properties(placeholder_type):
    """Create explicit theme references when a placeholder has no local rPr."""
    title_types = {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
    }
    family = 'mj' if placeholder_type in title_types else 'mn'
    run_props = etree.Element(_a('rPr'))
    run_props.set('lang', 'ko-KR')
    run_props.set('dirty', '0')
    for tag, script in (('latin', 'lt'), ('ea', 'ea'), ('cs', 'cs')):
        font = etree.SubElement(run_props, _a(tag))
        font.set('typeface', f'+{family}-{script}')
    return run_props


def _master_character_properties(layout, layout_placeholder):
    """Find character properties inherited from the corresponding master placeholder."""
    try:
        placeholder_type = layout_placeholder.placeholder_format.type
        candidates = [
            placeholder
            for placeholder in layout.slide_master.placeholders
            if placeholder.placeholder_format.type == placeholder_type
        ]
        for placeholder in candidates:
            props = _character_properties_from_text_body(placeholder.text_frame._txBody)
            if props is not None:
                return props
    except Exception:
        pass
    return None


def parse_lyrics_text(raw_text):
    lyrics_dict = {}
    blocks = re.split(r'\n\s*\n', raw_text.strip())

    for block in blocks:
        lines = block.strip().split('\n')
        if not lines:
            continue

        part_key = lines[0].strip().lower()
        if not part_key:
            continue
        lyrics_content = '\n'.join(line.strip() for line in lines[1:])
        lyrics_dict[part_key] = lyrics_content

    return lyrics_dict


def get_base_key(part_key):
    return re.sub(r"'+$", "", part_key)


def chunk_text(text, max_lines=2):
    if not text:
        return []
    lines = [line.strip() for line in text.split('\n') if line.strip()]

    chunks = []
    i = 0
    while i < len(lines):
        take_lines = max_lines
        if max_lines == 1 and len(lines[i]) <= _MIN_SHORT_LINE_LEN and i + 1 < len(lines):
            take_lines = 2

        chunk_lines = lines[i:i + take_lines]
        chunks.append('\n'.join(chunk_lines))
        i += len(chunk_lines)

    return chunks


def wrap_text_by_max_chars(text, max_chars_per_line=9999):
    if not text:
        return ""

    try:
        max_chars = int(max_chars_per_line)
    except (TypeError, ValueError):
        max_chars = 18

    if max_chars <= 0:
        return text

    wrapped_lines = []
    for raw_line in text.split("\n"):
        line = raw_line.strip()
        if not line:
            continue

        while len(line) > max_chars:
            break_at = max_chars
            space_at = line.rfind(" ", 0, max_chars + 1)
            if space_at >= max(1, max_chars // 2):
                break_at = space_at

            wrapped_lines.append(line[:break_at].rstrip())
            line = line[break_at:].lstrip()

        if line:
            wrapped_lines.append(line)

    return "\n".join(wrapped_lines)


def _read_layout_ph_fmt(shape):
    """슬라이드 placeholder에 대응하는 레이아웃 placeholder에서 lstStyle/pPr/rPr를 읽는다.
    슬라이드 placeholder는 add_slide() 직후 항상 빈 XML이므로 레이아웃에서 읽어야 한다."""
    try:
        ph_idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
        for lph in layout.placeholders:
            if lph.placeholder_format.idx != ph_idx:
                continue
            txb = lph.text_frame._txBody
            saved_lstStyle = None
            saved_pPr = None
            saved_rPr = _character_properties_from_text_body(txb)
            lst = txb.find(_a('lstStyle'))
            if lst is not None and len(lst):
                saved_lstStyle = deepcopy(lst)
            ps = txb.findall(f'{{{_NS}}}p')
            if ps:
                p0 = ps[0]
                pPr = p0.find(_a('pPr'))
                if pPr is not None:
                    saved_pPr = deepcopy(pPr)
            if saved_rPr is None:
                saved_rPr = _master_character_properties(layout, lph)
            if saved_rPr is None:
                saved_rPr = _default_theme_run_properties(lph.placeholder_format.type)
            return saved_lstStyle, saved_pPr, saved_rPr
    except Exception:
        pass
    return None, None, None


def set_editable_text(shape, text, font_size=None):
    """Set text while preserving the layout placeholder's font/paragraph formatting."""
    from lxml import etree as _etree
    tf = shape.text_frame
    txBody = tf._txBody

    saved_lstStyle, saved_pPr, saved_rPr = _read_layout_ph_fmt(shape)

    # lstStyle 교체: 슬라이드 placeholder의 빈 lstStyle을 레이아웃 것으로 채운다
    if saved_lstStyle is not None:
        existing = txBody.find(f'{{{_NS}}}lstStyle')
        if existing is not None:
            txBody.remove(existing)
        bodyPr = txBody.find(f'{{{_NS}}}bodyPr')
        insert_at = (list(txBody).index(bodyPr) + 1) if bodyPr is not None else 0
        txBody.insert(insert_at, deepcopy(saved_lstStyle))

    tf.clear()

    lines = text.split('\n') if text else [""]
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_elm = para._p

        if saved_pPr is not None:
            existing_pPr = p_elm.find(f'{{{_NS}}}pPr')
            if existing_pPr is not None:
                p_elm.remove(existing_pPr)
            p_elm.insert(0, deepcopy(saved_pPr))

        run = para.add_run()
        run.text = line
        r_elm = run._r

        if saved_rPr is not None:
            existing_rPr = r_elm.find(f'{{{_NS}}}rPr')
            if existing_rPr is not None:
                r_elm.remove(existing_rPr)
            new_rPr = _ensure_script_fonts(deepcopy(saved_rPr))
            if font_size:
                new_rPr.set('sz', str(int(font_size * 100)))
            r_elm.insert(0, new_rPr)
        else:
            # saved_rPr 없음: auto-generated rPr 재사용
            new_rPr = r_elm.find(f'{{{_NS}}}rPr')
            if new_rPr is None:
                new_rPr = _etree.Element(f'{{{_NS}}}rPr')
                new_rPr.set('lang', 'ko-KR')
                new_rPr.set('dirty', '0')
                r_elm.insert(0, new_rPr)
            if font_size:
                new_rPr.set('sz', str(int(font_size * 100)))

        # New typing at the end of a paragraph inherits endParaRPr. Keep it
        # aligned with the run so PowerPoint does not switch fonts after save.
        existing_end_props = p_elm.find(_a('endParaRPr'))
        if existing_end_props is not None:
            p_elm.remove(existing_end_props)
        end_props = deepcopy(new_rPr)
        end_props.tag = _a('endParaRPr')
        p_elm.append(end_props)


def delete_all_slides(prs):
    slide_id_list = prs.slides._sldIdLst
    for slide_id in tuple(slide_id_list):
        rel_id = slide_id.rId
        slide_id_list.remove(slide_id)
        prs.part.drop_rel(rel_id)

    if len(prs.slides) != 0:
        raise RuntimeError("템플릿의 기존 슬라이드를 삭제하지 못했습니다.")


def find_layout(prs, keywords, fallback=None):
    for layout in prs.slide_layouts:
        layout_name = layout.name.lower()
        if any(keyword.lower() in layout_name for keyword in keywords):
            return layout
    return fallback


def add_layout_slide(prs, keywords, fallback_keywords=None):
    fallback = find_layout(prs, fallback_keywords or []) if fallback_keywords else None
    layout = find_layout(prs, keywords, fallback)
    if layout is None:
        layout = prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def add_opening_slides(prs):
    """Add the fixed opening slides for an integrated lyrics PPT."""
    add_layout_slide(prs, ["home", "홈"])
    add_layout_slide(prs, ["예배를 시작하며", "worship"], fallback_keywords=["home", "홈"])


def reset_integrated_ppt(prs):
    """Remove all template slides so the output contains only generated slides."""
    delete_all_slides(prs)
    add_opening_slides(prs)


def append_closing_slide(prs):
    add_layout_slide(prs, ["기도", "prayer"], fallback_keywords=["home", "홈"])


def parse_sequence_text(sequence_text):
    lines = [line.strip() for line in sequence_text.splitlines() if line.strip()]

    if not lines:
        raise ValueError("레파토리 입력창이 비어 있습니다.")

    if len(lines) % 2 != 0:
        raise ValueError(f"'{lines[-1]}' 다음에 진행 순서 줄이 없습니다.")

    sequence_entries = []
    for i in range(0, len(lines), 2):
        song_title = lines[i]
        sequence_str = lines[i + 1]

        if not song_title:
            raise ValueError(f"{i + 1}번째 줄의 곡 제목이 비어 있습니다.")
        if not sequence_str:
            raise ValueError(f"'{song_title}'의 진행 순서가 비어 있습니다.")

        sequence_entries.append((song_title, sequence_str))

    return sequence_entries


def append_lyrics_to_ppt(
    prs,
    song_title,
    lyrics_text,
    sequence_str,
    max_lines_per_slide=2,
    max_chars_per_line=9999,
    lyrics_font_size=None,
):
    lyrics_dict = parse_lyrics_text(lyrics_text)
    sequence_list = [part.strip() for part in sequence_str.split('-') if part.strip()]

    title_layout = None
    lyrics_layout = None

    for layout in prs.slide_layouts:
        if "제목" in layout.name:
            title_layout = layout
        if "가사" in layout.name:
            lyrics_layout = layout

    if title_layout is None:
        title_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    if lyrics_layout is None:
        lyrics_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[0]

    title_slide = prs.slides.add_slide(title_layout)
    for shape in title_slide.placeholders:
        if shape.has_text_frame:
            set_editable_text(shape, song_title)
            break

    for idx, part in enumerate(sequence_list):
        part_lower = part.lower()
        base_part_lower = get_base_key(part_lower)

        if part_lower in lyrics_dict:
            display_text = lyrics_dict[part_lower]
        elif base_part_lower in lyrics_dict:
            display_text = lyrics_dict[base_part_lower]
        else:
            if part.startswith('(') and part.endswith(')'):
                display_text = part[1:-1].strip()
            else:
                if idx == 0 and base_part_lower in ("i", "intro"):
                    continue
                display_text = "-"

        display_text = wrap_text_by_max_chars(display_text, max_chars_per_line)
        chunks = chunk_text(display_text, max_lines_per_slide)
        if not chunks:
            chunks = [""]

        for chunk in chunks:
            slide = prs.slides.add_slide(lyrics_layout)

            placeholders = [shape for shape in slide.placeholders if shape.has_text_frame]
            placeholders.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)

            lyrics_placeholder = placeholders[0] if len(placeholders) > 0 else None
            song_title_placeholder = placeholders[1] if len(placeholders) > 1 else None

            if lyrics_placeholder is not None:
                set_editable_text(lyrics_placeholder, chunk, font_size=lyrics_font_size)

            if song_title_placeholder is not None:
                set_editable_text(song_title_placeholder, song_title)
