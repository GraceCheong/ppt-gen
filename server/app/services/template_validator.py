"""템플릿 PPTX 호환성 검증."""
from __future__ import annotations

from dataclasses import dataclass, field


@dataclass
class TemplateValidationResult:
    compatible: bool
    issues: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    layout_names: list[str] = field(default_factory=list)


def validate_template(pptx_path: str) -> TemplateValidationResult:
    """ppt_builder.py의 슬라이드 생성 로직 기준으로 호환성을 검사한다.

    필수 (없으면 compatible=False):
    - "가사" 이름을 포함한 레이아웃 + 텍스트 placeholder 2개 이상
    - "제목" 이름을 포함한 레이아웃

    권장 (없어도 동작하나 경고):
    - "홈" 또는 "home" 레이아웃 (fallback이 slide_layouts[0])
    - "예배를 시작하며" 또는 "worship" 레이아웃
    - "기도" 또는 "prayer" 레이아웃
    """
    from pptx import Presentation

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return TemplateValidationResult(
            compatible=False,
            issues=[f"파일을 열 수 없습니다: {e}"],
        )

    layout_names = [layout.name for layout in prs.slide_layouts]
    lower_names = [n.lower() for n in layout_names]

    issues: list[str] = []
    warnings: list[str] = []

    # ── 필수: "가사" 레이아웃 ──────────────────────────────────────────────
    lyrics_layout = next(
        (layout for layout in prs.slide_layouts if "가사" in layout.name),
        None,
    )
    if lyrics_layout is None:
        issues.append("lyrics_layout_missing")
    else:
        text_phs = [ph for ph in lyrics_layout.placeholders if ph.has_text_frame]
        if len(text_phs) < 2:
            issues.append("lyrics_placeholder_count")

    # ── 필수: "제목" 레이아웃 ────────────────────────────────────────────────
    has_title = any("제목" in name for name in layout_names)
    if not has_title:
        issues.append("title_layout_missing")

    # ── 권장: "홈" 레이아웃 ──────────────────────────────────────────────────
    has_home = any("홈" in n or "home" in n for n in lower_names)
    if not has_home:
        warnings.append("home_layout_missing")

    # ── 권장: "예배를 시작하며" 또는 "worship" 레이아웃 ─────────────────────
    has_worship = any("예배를 시작하며" in n or "worship" in n for n in lower_names)
    if not has_worship:
        warnings.append("worship_layout_missing")

    # ── 권장: "기도" 또는 "prayer" 레이아웃 ──────────────────────────────────
    has_prayer = any("기도" in n or "prayer" in n for n in lower_names)
    if not has_prayer:
        warnings.append("prayer_layout_missing")

    return TemplateValidationResult(
        compatible=len(issues) == 0,
        issues=issues,
        warnings=warnings,
        layout_names=layout_names,
    )
