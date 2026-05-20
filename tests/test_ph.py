from pathlib import Path

from pptx import Presentation


ROOT_DIR = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = ROOT_DIR / "assets" / "templates" / "template 2.pptx"


def test_template_has_expected_layouts():
    prs = Presentation(str(TEMPLATE_PATH))
    layout_names = [layout.name for layout in prs.slide_layouts]

    assert any("제목" in name or "Title" in name for name in layout_names)
    assert any("가사" in name or "Lyrics" in name for name in layout_names)


def test_lyrics_layout_has_text_placeholder():
    prs = Presentation(str(TEMPLATE_PATH))
    lyrics_layout = next(
        layout
        for layout in prs.slide_layouts
        if "가사" in layout.name or "Lyrics" in layout.name
    )

    text_placeholders = [
        placeholder
        for placeholder in lyrics_layout.placeholders
        if placeholder.has_text_frame
    ]

    assert text_placeholders
