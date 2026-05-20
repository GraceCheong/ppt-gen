from pathlib import Path
import sys

from pptx import Presentation


ROOT_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT_DIR / "src"))

from ppt_builder import append_lyrics_to_ppt, chunk_text, parse_lyrics_text, wrap_text_by_max_chars
from ppt_service import build_integrated_pptx


TEMPLATE_PATH = ROOT_DIR / "assets" / "templates" / "template 2.pptx"


SAMPLE_LYRICS = """V1
가사 첫 번째 줄
가사 두 번째 줄

C
후렴 첫 번째 줄
후렴 두 번째 줄

Out
마지막 가사"""


def test_parse_lyrics_text():
    parsed = parse_lyrics_text(SAMPLE_LYRICS)

    assert parsed["V1"] == "가사 첫 번째 줄\n가사 두 번째 줄"
    assert parsed["C"] == "후렴 첫 번째 줄\n후렴 두 번째 줄"
    assert parsed["Out"] == "마지막 가사"


def test_chunk_text_respects_max_lines():
    chunks = chunk_text("첫 줄\n둘째 줄\n셋째 줄", max_lines=2)

    assert chunks == ["첫 줄\n둘째 줄", "셋째 줄"]


def test_wrap_text_by_max_chars_splits_long_lines():
    wrapped = wrap_text_by_max_chars("1234567890", max_chars_per_line=4)

    assert wrapped == "1234\n5678\n90"


def test_append_lyrics_to_ppt_adds_slides():
    prs = Presentation(str(TEMPLATE_PATH))
    before_count = len(prs.slides)

    append_lyrics_to_ppt(
        prs,
        "테스트 곡",
        SAMPLE_LYRICS,
        "I-V1-C-Out",
        max_lines_per_slide=2,
    )

    assert len(prs.slides) > before_count


def test_append_lyrics_to_ppt_keeps_user_selected_four_line_chunks():
    prs = Presentation(str(TEMPLATE_PATH))
    before_count = len(prs.slides)
    long_lyrics = """V
이 줄은 긴 줄 기준보다 충분히 긴 가사입니다
둘째 줄
셋째 줄
넷째 줄"""

    append_lyrics_to_ppt(
        prs,
        "네 줄 테스트",
        long_lyrics,
        "V",
        max_lines_per_slide=4,
        max_chars_per_line=100,
    )

    assert len(prs.slides) == before_count + 2


def test_append_lyrics_to_ppt_applies_lyrics_font_size():
    prs = Presentation(str(TEMPLATE_PATH))

    append_lyrics_to_ppt(
        prs,
        "폰트 테스트",
        "V\n가사 첫 줄",
        "V",
        max_lines_per_slide=4,
        lyrics_font_size=36,
    )

    lyrics_slide = prs.slides[-1]
    placeholders = [shape for shape in lyrics_slide.placeholders if shape.has_text_frame]
    placeholders.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
    run = placeholders[0].text_frame.paragraphs[0].runs[0]

    assert run.font.size.pt == 36


def test_build_integrated_pptx_ignores_template_slides_and_adds_expected_shell(tmp_path):
    output_path = tmp_path / "out.pptx"

    result = build_integrated_pptx(
        str(TEMPLATE_PATH),
        [("테스트 곡", "V1")],
        {"테스트 곡": "V1\n첫 줄\n둘째 줄"},
        str(output_path),
        max_lines_per_slide=4,
    )

    assert result["appended_count"] == 1

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 5

    layout_names = [slide.slide_layout.name for slide in prs.slides]
    assert layout_names[0] == "Home"
    assert "예배를 시작하며" in layout_names[1]
    assert "제목" in layout_names[2]
    assert "가사" in layout_names[3]
    assert layout_names[4] == "Home"
