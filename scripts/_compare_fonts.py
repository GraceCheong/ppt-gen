"""
두 폰트(Diphylleia vs 맑은 고딕)로 각각 렌더링해서 비교 이미지 생성
어느 폰트가 현재 렌더링되고 있는지 시각적으로 확인
"""
import sys, os, tempfile, subprocess, copy, zipfile
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(__file__)), "src"))

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from songlist_builder import build_songlist_pptx, find_libreoffice, _patch_theme_korean_font, _TITLE_FONT

TEMPLATE = os.path.join(os.path.dirname(os.path.dirname(__file__)), "assets", "templates", "songlist_template.pptx")
TITLES = ["믿음이 없이는", "주님 마음 내게 주소서", "주의 자녀로 산다는 것은"]
OUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "out")
os.makedirs(OUT_DIR, exist_ok=True)

lo = find_libreoffice()
if not lo:
    print("LibreOffice 없음"); sys.exit(1)


def render_with_font(font_name, out_prefix):
    """지정 폰트로 생성 후 LibreOffice로 렌더링"""
    prs = Presentation(TEMPLATE)

    # theme에서 Hang 폰트 강제 설정
    _THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    for master in prs.slide_masters:
        try:
            theme_part = master.part.part_related_by(_THEME_REL)
        except Exception:
            continue
        root = etree.fromstring(theme_part.blob)
        fs = root.find(".//" + qn("a:fontScheme"))
        if fs is None:
            continue
        for section_tag in ("a:majorFont", "a:minorFont"):
            section = fs.find(qn(section_tag))
            if section is None:
                continue
            ea = section.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(section, qn("a:ea"))
            ea.set("typeface", font_name)
            for fe in section.findall(qn("a:font")):
                if fe.get("script") == "Hang":
                    fe.set("typeface", font_name)
                    break
        theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            alt = shape._element.nvSpPr.cNvPr.get("descr", "") or ""
        except AttributeError:
            alt = ""
        if alt.startswith("곡") and alt[1:].strip().isdigit():
            idx = int(alt[1:].strip()) - 1
            title = TITLES[idx] if idx < len(TITLES) else ""
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = title
            endParaRPr = para._p.find(qn("a:endParaRPr"))
            if endParaRPr is not None:
                rPr = copy.deepcopy(endParaRPr)
                rPr.tag = qn("a:rPr")
                for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
                    elem = rPr.find(qn(tag))
                    if elem is None:
                        elem = etree.SubElement(rPr, qn(tag))
                    elem.set("typeface", font_name)
                ex = run._r.find(qn("a:rPr"))
                if ex is not None:
                    run._r.remove(ex)
                run._r.insert(0, rPr)

    pptx_path = os.path.join(OUT_DIR, f"{out_prefix}.pptx")
    png_path  = os.path.join(OUT_DIR, f"{out_prefix}.png")
    prs.save(pptx_path)

    subprocess.run(
        [lo, "--headless", "--norestore", "--convert-to", "png", "--outdir", OUT_DIR, pptx_path],
        capture_output=True, timeout=60
    )
    return png_path


fonts = [
    ("Diphylleia",   "compare_diphylleia"),
    ("맑은 고딕",     "compare_malgun"),
    ("나눔명조",      "compare_nanum"),
]

for font_name, prefix in fonts:
    path = render_with_font(font_name, prefix)
    exists = os.path.exists(path)
    print(f"{font_name:15s} → {path}  {'✓' if exists else '✗'}")

print("\n각 PNG를 열어서 현재 렌더링과 비교하세요.")
