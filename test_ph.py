import os
from pptx import Presentation

prs = Presentation('template.pptx')
l = prs.slide_layouts[3]
print("Layout: ", l.name)
for idx, p in enumerate(l.placeholders):
    print(f"[{idx}] name={p.name}, has_text_frame={p.has_text_frame}, width={p.width}, height={p.height}")
